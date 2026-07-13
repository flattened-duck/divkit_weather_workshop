"""deck.py — on-brand slide-building helpers for the YaTeam corporate template.

Loads the corporate template (keeps all masters / layouts / theme), strips the
184 demo slides, and exposes friendly builders so the build agent focuses on
layout, not python-pptx boilerplate.

Usage:
    import deck
    prs = deck.new_deck()
    s = deck.add(prs, "title")
    deck.set_title(s, "BDUI и DivKit")
    deck.fill_bodies(s, ["Опыт большого проекта", "Имя · должность"])
    ...
    prs.save("/path/out.pptx")

Then render to PNG per slide with render.py.

Coordinate system: centimetres. Slide is 67.74 x 38.1 cm (16:9).
Content safe-zone the template uses: left margin ~3.8 cm, top title ~2.45 cm.
"""

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TEMPLATE = "/Users/the-leo/Downloads/Shablon_ЯTeam_Power Point_светлая_красно-желтые элементы.pptx"

# ---- Brand palette (from the template theme) -------------------------------
BRAND = {
    "ink":        "222A3A",  # primary text (dk1)
    "ink2":       "575E6E",  # secondary text (dk2)
    "white":      "FFFFFF",
    "lilac_bg":   "F5F4FD",  # light surface (lt2)
    "red":        "F8604A",  # accent1 — primary brand accent
    "yellow":     "FFC812",  # accent2
    "pink_bg":    "FFF4F2",  # accent3 — soft red surface
    "yellow_bg":  "FFF9E6",  # accent4 — soft yellow surface
    "grey":       "CFCDDD",  # accent5
}
HEAD_FONT = "YS Text Medium"
BODY_FONT = "YS Text Regular"

SLIDE_W_CM = 67.74
SLIDE_H_CM = 38.1
MARGIN_L = 3.8

# Friendly layout name -> exact template layout name (resolved across masters).
LAYOUTS = {
    "title":               "Титульный 1",
    "title_plain":         "1_Титульный 1",
    "agenda":              "Содержание 1",
    "header_text":         "Заголовок в две строки + текст",
    "title_subtitle":      "Заголовок + подзаголовок",
    "title_subtitle_text": "Заголовок + подзаголовок + текст",
    "three_blocks":        "3 блока текста 1",
    "two_blocks":          "2 текстовых блока",
    "compare":             "Сравнение",
    "quote":               "Важная мысль, цитата 1",
    "big_thesis":          "Заголлвок + большой тезис",  # sic: template typo
    "numbered":            "Нумерованный список",
    "stages":              "Этапы 1",
    "cards":               "карточки 2",
    "blank":               "Пустой",
    "questions":           "Вопросы + QR",
}


def rgb(name_or_hex):
    return RGBColor.from_string(BRAND.get(name_or_hex, name_or_hex))


# ---- Deck lifecycle --------------------------------------------------------
def new_deck():
    """Load template, remove all demo slides, return an empty on-brand deck."""
    prs = Presentation(TEMPLATE)
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        rId = sld_id.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sld_id_lst.remove(sld_id)
    return prs


def _resolve_layout(prs, friendly):
    """Find a layout object by friendly key (or raw exact name) across masters."""
    target = LAYOUTS.get(friendly, friendly)
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == target:
                return layout
    raise KeyError(f"layout not found: {friendly!r} -> {target!r}")


def add(prs, friendly="blank"):
    """Add a slide using a friendly layout name. Returns the slide."""
    layout = _resolve_layout(prs, friendly)
    return prs.slides.add_slide(layout)


# ---- Placeholder fillers ---------------------------------------------------
def _placeholders_by_type(slide):
    from pptx.enum.shapes import PP_PLACEHOLDER
    title, bodies = None, []
    for ph in slide.placeholders:
        t = ph.placeholder_format.type
        if t == PP_PLACEHOLDER.TITLE:
            title = ph
        elif t == PP_PLACEHOLDER.BODY:
            bodies.append(ph)
    return title, bodies


def set_title(slide, text):
    title, _ = _placeholders_by_type(slide)
    if title is None:
        return None
    title.text_frame.text = text
    return title


def fill_bodies(slide, texts):
    """Fill BODY placeholders in document order. `texts` is a list; an item may
    be a string or a list of strings (multi-paragraph / bullets)."""
    _, bodies = _placeholders_by_type(slide)
    for ph, content in zip(bodies, texts):
        tf = ph.text_frame
        if isinstance(content, (list, tuple)):
            tf.text = str(content[0])
            for line in content[1:]:
                p = tf.add_paragraph()
                p.text = str(line)
        else:
            tf.text = str(content)
    return bodies


# ---- Free-form primitives (for blank-canvas diagrams) ----------------------
def textbox(slide, x, y, w, h, text, size=18, color="ink", bold=False,
            font=None, align="left", anchor="top", line_spacing=None):
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    lines = text.split("\n") if isinstance(text, str) else list(text)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font or (HEAD_FONT if bold else BODY_FONT)
        run.font.color.rgb = rgb(color)
    return tb


def card(slide, x, y, w, h, fill="lilac_bg", radius=0.12, line=None):
    """Rounded-rectangle surface. Returns the shape (add text via .text_frame
    or use textbox on top)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    try:  # adjust corner radius
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def card_with_text(slide, x, y, w, h, title=None, body=None,
                   fill="lilac_bg", accent="red", title_size=20, body_size=15,
                   pad=0.8):
    card(slide, x, y, w, h, fill=fill)
    ty = y + pad
    if title:
        textbox(slide, x + pad, ty, w - 2 * pad, 2.0, title,
                size=title_size, bold=True, color="ink")
        ty += 2.0
    if body:
        textbox(slide, x + pad, ty, w - 2 * pad, h - (ty - y) - pad, body,
                size=body_size, color="ink2")


def chip(slide, x, y, text, fill="red", color="white", size=13, w=None, h=1.3):
    if w is None:
        w = 0.7 + 0.42 * len(text)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = HEAD_FONT
    r.font.color.rgb = rgb(color)
    return shp


def arrow(slide, x1, y1, x2, y2, color="ink2", weight=2.5):
    """Straight connector with an arrow head from (x1,y1) to (x2,y2) in cm."""
    conn = slide.shapes.add_connector(2, Cm(x1), Cm(y1), Cm(x2), Cm(y2))  # 2=STRAIGHT
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(weight)
    # add arrowhead on the end
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def divider(slide, x, y, w, color="red", weight=3.0):
    return arrow_line(slide, x, y, x + w, y, color=color, weight=weight, head=False)


def arrow_line(slide, x1, y1, x2, y2, color="ink2", weight=2.5, head=True):
    conn = slide.shapes.add_connector(2, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(weight)
    if head:
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"),
                              {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return conn


def number_badge(slide, x, y, n, d=1.6, fill="red", color="white", size=20):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(x), Cm(y), Cm(d), Cm(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    shp.shadow.inherit = False
    tf = shp.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = str(n)
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.name = HEAD_FONT
    r.font.color.rgb = rgb(color)
    return shp


def notes(slide, text):
    """Set speaker notes (used to embed lecture talking points)."""
    slide.notes_slide.notes_text_frame.text = text


def image(slide, path, x, y, w, h):
    """Insert a real image (for later, when screenshots exist)."""
    return slide.shapes.add_picture(path, Cm(x), Cm(y), Cm(w), Cm(h))


def placeholder(slide, x, y, w, h, kind="screen", caption=""):
    """Framed dashed placeholder box for an asset to be added later.
    kind: 'screen' | 'code' | 'screencast'."""
    icon = {"screen": "🖼  СКРИН", "code": "</>  КОД",
            "screencast": "▶  СКРИНКАСТ"}.get(kind, "□")
    fill = {"screen": "lilac_bg", "code": "yellow_bg",
            "screencast": "pink_bg"}.get(kind, "lilac_bg")
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Cm(x), Cm(y), Cm(w), Cm(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb("grey")
    shp.line.width = Pt(1.5)
    ln = shp.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.04
    except Exception:
        pass
    label = icon if not caption else f"{icon}\n{caption}"
    textbox(slide, x, y, w, h, label, size=16, color="ink2", bold=True,
            align="center", anchor="middle")
    return shp
