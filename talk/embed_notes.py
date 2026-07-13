"""Embed per-slide talking points from lecture_notes.md into pptx speaker notes.

The corporate template's notes master has NO placeholders, so python-pptx can't
auto-create a notes text frame. We inject a body placeholder into each notes
slide via XML, then fill it.
"""
import re
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

SRC = "/Users/the-leo/divkit-weather-workshop/talk/bdui_divkit.pptx"
NOTES_MD = "/Users/the-leo/divkit-weather-workshop/talk/lecture_notes.md"
OUT = "/Users/the-leo/divkit-weather-workshop/talk/BDUI_DivKit_лекция.pptx"

# Standard notes body geometry
OFF_X, OFF_Y, EXT_CX, EXT_CY = 685800, 2971800, 5486400, 3884613


def parse_notes(md):
    parts = re.split(r"(?m)^##\s+Слайд\s+(\d+)\s*—\s*(.*)$", md)
    out = {}
    for i in range(1, len(parts), 3):
        n = int(parts[i])
        title = parts[i + 1].strip()
        bullets = [ln.strip()[2:].strip() for ln in parts[i + 2].splitlines()
                   if ln.strip().startswith("- ")]
        out[n] = (title, bullets)
    return out


def ensure_notes_tf(notes_slide):
    """Return a text frame for the notes body, creating a body placeholder if
    the notes slide has none."""
    tf = notes_slide.notes_text_frame
    if tf is not None:
        return tf
    spTree = notes_slide.shapes._spTree
    a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    p = "http://schemas.openxmlformats.org/presentationml/2006/main"
    xml = (
        f'<p:sp xmlns:p="{p}" xmlns:a="{a}">'
        '<p:nvSpPr>'
        '<p:cNvPr id="2" name="Notes Placeholder"/>'
        '<p:cNvSpPr><a:spLocks noGrp="1"/></p:cNvSpPr>'
        '<p:nvPr><p:ph type="body" idx="1"/></p:nvPr>'
        '</p:nvSpPr>'
        '<p:spPr>'
        f'<a:xfrm><a:off x="{OFF_X}" y="{OFF_Y}"/><a:ext cx="{EXT_CX}" cy="{EXT_CY}"/></a:xfrm>'
        '</p:spPr>'
        '<p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>'
        '</p:sp>'
    )
    from pptx.oxml import parse_xml
    spTree.append(parse_xml(xml))
    return notes_slide.notes_text_frame


def main():
    notes = parse_notes(open(NOTES_MD, encoding="utf-8").read())
    prs = Presentation(SRC)
    total = len(prs.slides._sldIdLst)
    applied = 0
    for idx, slide in enumerate(prs.slides, start=1):
        if idx not in notes:
            continue
        title, bullets = notes[idx]
        tf = ensure_notes_tf(slide.notes_slide)
        tf.text = f"Слайд {idx} — {title}"
        for b in bullets:
            par = tf.add_paragraph()
            par.text = "• " + b
        applied += 1
    prs.save(OUT)
    print(f"slides={total}  notes_applied={applied}  parsed={len(notes)}")

    # verify
    prs2 = Presentation(OUT)
    s5 = list(prs2.slides)[4]
    print("--- slide 5 notes ---")
    print(s5.notes_slide.notes_text_frame.text[:280])


if __name__ == "__main__":
    main()
