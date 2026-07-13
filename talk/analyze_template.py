"""Dump the structure of the corporate pptx template: slide size, theme fonts,
theme colors, and every slide layout with its placeholders + the design slides
already present in the file. Read-only."""

import sys
from pptx import Presentation
from pptx.util import Emu

TEMPLATE = "/Users/the-leo/Downloads/Shablon_ЯTeam_Power Point_светлая_красно-желтые элементы.pptx"


def emu_to_cm(v):
    if v is None:
        return None
    return round(Emu(v).cm, 2)


def main():
    prs = Presentation(TEMPLATE)
    print("=== SLIDE SIZE ===")
    print(f"width  = {emu_to_cm(prs.slide_width)} cm")
    print(f"height = {emu_to_cm(prs.slide_height)} cm")
    aspect = prs.slide_width / prs.slide_height
    print(f"aspect = {round(aspect, 3)}  ({'16:9' if abs(aspect-16/9)<0.02 else '4:3' if abs(aspect-4/3)<0.02 else 'custom'})")

    # Theme fonts + colors via the first master's theme part XML
    print("\n=== THEME (fonts + color scheme) ===")
    try:
        master = prs.slide_masters[0]
        theme = master.element.getroottree()  # placeholder; real theme below
    except Exception as e:
        print("theme introspection error:", e)

    # python-pptx exposes theme via part; parse XML directly
    from pptx.oxml.ns import qn
    for mi, master in enumerate(prs.slide_masters):
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        import lxml.etree as ET
        root = ET.fromstring(theme_part.blob)
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clrs = root.find(".//a:clrScheme", ns)
        print(f"\n-- master[{mi}] color scheme: {clrs.get('name') if clrs is not None else '?'}")
        if clrs is not None:
            for child in clrs:
                tag = ET.QName(child).localname
                srgb = child.find("a:srgbClr", ns)
                sysv = child.find("a:sysClr", ns)
                val = srgb.get("val") if srgb is not None else (sysv.get("lastClr") if sysv is not None else "?")
                print(f"   {tag:14s} #{val}")
        fonts = root.find(".//a:fontScheme", ns)
        if fonts is not None:
            major = fonts.find("a:majorFont/a:latin", ns)
            minor = fonts.find("a:minorFont/a:latin", ns)
            print(f"   majorFont = {major.get('typeface') if major is not None else '?'}")
            print(f"   minorFont = {minor.get('typeface') if minor is not None else '?'}")

    print("\n=== SLIDE LAYOUTS ===")
    for mi, master in enumerate(prs.slide_masters):
        print(f"\n## master[{mi}] name={master.name!r} layouts={len(master.slide_layouts)}")
        for li, layout in enumerate(master.slide_layouts):
            phs = list(layout.placeholders)
            print(f"  layout[{li}] {layout.name!r}  placeholders={len(phs)}")
            for ph in phs:
                pf = ph.placeholder_format
                print(
                    f"      idx={pf.idx} type={pf.type} name={ph.name!r} "
                    f"pos=({emu_to_cm(ph.left)},{emu_to_cm(ph.top)}) "
                    f"size=({emu_to_cm(ph.width)}x{emu_to_cm(ph.height)})"
                )

    print("\n=== EXISTING DESIGN SLIDES ===")
    print(f"slide count in template = {len(prs.slides)}")
    for si, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        texts = []
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                t = sh.text_frame.text.strip().replace("\n", " | ")
                texts.append(t[:60])
        print(f"  slide[{si}] layout={layout_name!r}  text={texts[:3]}")


if __name__ == "__main__":
    main()
