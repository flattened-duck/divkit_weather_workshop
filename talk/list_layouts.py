"""Compact dump: every slide layout name (per master) + any existing design
slides in the template. Helps choose on-brand layouts when building."""

from pptx import Presentation

TEMPLATE = "/Users/the-leo/Downloads/Shablon_ЯTeam_Power Point_светлая_красно-желтые элементы.pptx"

prs = Presentation(TEMPLATE)

for mi, master in enumerate(prs.slide_masters):
    print(f"\n===== MASTER[{mi}] layouts={len(master.slide_layouts)} =====")
    for li, layout in enumerate(master.slide_layouts):
        # count placeholder types for a quick sense of the layout
        kinds = {}
        for ph in layout.placeholders:
            k = str(ph.placeholder_format.type)
            kinds[k] = kinds.get(k, 0) + 1
        kinds_s = ",".join(f"{k}:{v}" for k, v in sorted(kinds.items()))
        print(f"  [{mi}.{li:>3}] {layout.name!r:50s} {kinds_s}")

print(f"\n===== EXISTING DESIGN SLIDES = {len(prs.slides)} =====")
for si, slide in enumerate(prs.slides):
    txt = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt.append(sh.text_frame.text.strip().replace("\n", " | ")[:50])
    print(f"  slide[{si}] layout={slide.slide_layout.name!r} :: {txt[:2]}")
